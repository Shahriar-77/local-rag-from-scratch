
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Dict, List
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
import pandas as pd

class BaseIngestor(ABC):
    def __init__(self, file_path: Path):
        self.file_path = file_path

    @abstractmethod
    def extract(self) -> Dict:
        """
        Returns a dictionary with:
        - source_path
        - file_type
        - blocks      (structured content units)
        - text        (flattened debug view)
        - metadata
        """
        pass

class PDFIngestor(BaseIngestor):
    def extract(self) -> Dict:
        reader = PdfReader(self.file_path)
        blocks: List[Dict] = []

        for page_idx, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                blocks.append({
                    "type": "page",
                    "page_index": page_idx,
                    "text": text
                })

        flattened_text = "\n\n".join(b["text"] for b in blocks)

        return {
            "source_path": str(self.file_path),
            "file_type": "pdf",
            "blocks": blocks,
            "text": flattened_text,
            "metadata": {
                "num_pages": len(reader.pages)
            }
        }


class DOCXIngestor(BaseIngestor):
    def extract(self) -> Dict:
        doc = Document(self.file_path)
        blocks: List[Dict] = []

        for p in doc.paragraphs:
            text = p.text.strip()
            if not text:
                continue

            style_name = p.style.name if p.style else "Normal"

            if style_name.startswith("Heading"):
                try:
                    level = int(style_name.split()[-1])
                except ValueError:
                    level = None

                blocks.append({
                    "type": "heading",
                    "level": level,
                    "text": text
                })
            else:
                blocks.append({
                    "type": "paragraph",
                    "text": text
                })

        flattened_text = "\n".join(b["text"] for b in blocks)

        return {
            "source_path": str(self.file_path),
            "file_type": "docx",
            "blocks": blocks,
            "text": flattened_text,
            "metadata": {
                "num_blocks": len(blocks)
            }
        }


class PPTXIngestor(BaseIngestor):
    def extract(self) -> Dict:
        prs = Presentation(self.file_path)
        blocks: List[Dict] = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_elements = []

            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue

                text = shape.text.strip()
                if not text:
                    continue

                role = "body"

                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    if ph_type == PP_PLACEHOLDER.TITLE:
                        role = "title"

                slide_elements.append({
                    "role": role,
                    "text": text
                })

            if slide_elements:
                blocks.append({
                    "type": "slide",
                    "slide_index": slide_idx,
                    "elements": slide_elements
                })

        flattened_text = "\n\n".join(
            el["text"]
            for slide in blocks
            for el in slide["elements"]
        )

        return {
            "source_path": str(self.file_path),
            "file_type": "pptx",
            "blocks": blocks,
            "text": flattened_text,
            "metadata": {
                "num_slides": len(prs.slides)
            }
        }

class CSVIngestor(BaseIngestor):
    def extract(self) -> Dict:
        df = pd.read_csv(self.file_path)

        schema = {}
        for col in df.columns:
            series = df[col]

            inferred_type = str(series.dtype)

            # crude but effective datetime detection
            if inferred_type == "object":
                sample = series.dropna().iloc[:5]

                if not sample.empty:
                    try:
                        pd.to_datetime(sample, errors="raise")
                        inferred_type = "datetime"
                    except Exception:
                        pass



            schema[col] = inferred_type

        blocks = [{
            "type": "table",
            "rows": df.to_dict(orient="records")
        }]

        flattened_text = df.to_csv(index=False)

        return {
            "source_path": str(self.file_path),
            "file_type": "csv",
            "blocks": blocks,
            "text": flattened_text,
            "metadata": {
                "num_rows": df.shape[0],
                "num_columns": df.shape[1],
                "schema": schema
            }
        }
    

class TXTIngestor(BaseIngestor):
    def extract(self) -> Dict:
        content = self.file_path.read_text(encoding="utf-8")

        blocks = [{
            "type": "text",
            "text": content
        }]

        return {
            "source_path": str(self.file_path),
            "file_type": "txt",
            "blocks": blocks,
            "text": content,
            "metadata": {
                "num_characters": len(content)
            }
        }


class XLSXIngestor(BaseIngestor):
    def extract(self) -> dict:
        xls = pd.ExcelFile(self.file_path)
        blocks = []

        for sheet_name in xls.sheet_names:
            df = xls.parse(sheet_name)

            blocks.append({
                "type": "sheet",
                "sheet_name": sheet_name,
                "dataframe": df
            })

        combined_text = "\n\n".join(
            f"Sheet: {b['sheet_name']}\n{b['dataframe'].to_csv(index=False)}"
            for b in blocks
        )

        return {
            "source_path": str(self.file_path),
            "file_type": "xlsx",
            "blocks": blocks,
            "text": combined_text,
            "metadata": {
                "num_sheets": len(blocks),
                "sheet_names": xls.sheet_names
            }
        }
    
def get_ingestor(file_path: Path):
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        return PDFIngestor(file_path)
    elif suffix == ".docx":
        return DOCXIngestor(file_path)
    elif suffix == ".pptx":
        return PPTXIngestor(file_path)
    elif suffix == ".csv":
        return CSVIngestor(file_path)
    elif suffix == ".txt":
        return TXTIngestor(file_path)
    elif suffix =='.xlsx':
        return XLSXIngestor(file_path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")
    

def ingest_file(path_str: str):
    path = Path(path_str)

    if not path.exists():
        raise FileNotFoundError(path)

    ingestor = get_ingestor(path)
    document = ingestor.extract()

    return document

